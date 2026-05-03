from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color palette
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT     = RGBColor(0x00, 0xC8, 0xFF)
ACCENT2    = RGBColor(0x7B, 0x2F, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE8)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
GREEN      = RGBColor(0x00, 0xE5, 0x96)
RED_SOFT   = RGBColor(0xFF, 0x55, 0x55)
DARK_CARD  = RGBColor(0x0C, 0x1E, 0x35)
DARK_CARD2 = RGBColor(0x10, 0x22, 0x38)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─── Helpers ────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_DARK
    return s

def box(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def tb(slide, text, l, t, w, h, sz=16, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return txb

def header(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 0.06, fill=ACCENT)
    box(slide, 0, 0.06, 13.33, 1.15, fill=RGBColor(0x08,0x12,0x1E))
    tb(slide, title,    0.4, 0.12, 12.5, 0.68, sz=30, bold=True)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12.5, 0.38, sz=14,
           color=ACCENT, italic=True)

def pgnum(slide, n, total=7):
    tb(slide, f"{n} / {total}", 12.45, 7.1, 0.7, 0.32,
       sz=11, color=RGBColor(0x66,0x77,0x99), align=PP_ALIGN.RIGHT)

def two_run_line(slide, l, t, w, h, r1_text, r1_color, r1_sz, r1_bold,
                 r2_text, r2_color, r2_sz):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]
    def make_run(txt, col, sz, bold=False):
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = col; r.font.name = "Calibri"
    make_run(r1_text, r1_color, r1_sz, r1_bold)
    make_run(r2_text, r2_color, r2_sz)

# ─── SLIDE 1 — Title ────────────────────────────

s1 = new_slide()
box(s1, 0, 0,    13.33, 0.08, fill=ACCENT)
box(s1, 0, 7.42, 13.33, 0.08, fill=ACCENT2)
box(s1, 0, 0.08, 0.18, 7.34,  fill=ACCENT2)
box(s1, 9.0, 1.4, 4.0, 4.0,   fill=RGBColor(0x10,0x25,0x40))

tb(s1, "Autonomous Multi-Agent Swarms",
   0.5, 1.0, 12.0, 1.05, sz=38, bold=True)
tb(s1, "for Code Repository Remediation",
   0.5, 1.95, 11.0, 0.95, sz=38, bold=True, color=ACCENT)
tb(s1, "A Graph-Based Orchestration Approach",
   0.5, 3.0, 10.0, 0.5, sz=18, color=LIGHT_GRAY, italic=True)
box(s1, 0.5, 3.6, 5.5, 0.045, fill=ACCENT)

details = [
    ("Presenter:  ", "Durvankur Rasal"),
    ("Roll No:    ", "XXXXX"),
    ("Course:     ", "Advanced Software Engineering / AI Systems"),
    ("Date:       ", "April 17, 2026"),
]
y = 3.78
for lbl, val in details:
    two_run_line(s1, 0.5, y, 9.5, 0.38,
                 lbl, ACCENT, 14, True,
                 val, LIGHT_GRAY, 14)
    y += 0.42

tb(s1, "📌 Speaker Note:  Open with the elevator pitch — "
       "\"A Multi-Agent Swarm is a coordinated network of specialized AI agents "
       "that collaborate, debate, and self-correct to accomplish tasks no single model can handle well.\"",
   0.5, 6.48, 12.3, 0.85, sz=11, color=RGBColor(0x77,0x88,0xAA), italic=True)

pgnum(s1, 1)

# ─── SLIDE 2 — Problem ──────────────────────────

s2 = new_slide()
header(s2, "The Bottleneck in Modern Software Development", "Problem Definition")

cards = [
    (ACCENT2, "Context-Switching Penalty",
     "Every time a developer stops to review a large PR, it takes 20+ minutes to rebuild mental context. "
     "Across a whole team, this compounds into a massive hidden productivity cost."),
    (ACCENT,  "The 'Myopic AI' Flaw",
     "Standard LLMs see only the file diff — they have no idea how that file interacts with "
     "the rest of the codebase. The result is shallow reviews that miss systemic architectural issues."),
    (YELLOW,  "The Goal",
     "Shift from rigid, static CI/CD pipelines to a dynamic, context-aware remediation system "
     "that understands the full architecture before ever suggesting a fix."),
]

yc = 1.38
for bar_col, heading, body in cards:
    box(s2, 0.32, yc, 0.1, 1.52, fill=bar_col)
    box(s2, 0.46, yc, 12.5, 1.52, fill=DARK_CARD2)
    tb(s2, heading, 0.62, yc+0.1, 11.8, 0.42, sz=16, bold=True)
    tb(s2, body,    0.62, yc+0.52, 12.0, 0.9,  sz=13.5, color=LIGHT_GRAY)
    yc += 1.74

pgnum(s2, 2)

# ─── SLIDE 3 — Framework ────────────────────────

s3 = new_slide()
header(s3, "The Multi-Agent Framework", "From a single model to a team of specialized agents")

rows = [
    ("1", ACCENT2,
     "Division of Labor:  ",
     "One massive model that tries to be an expert in everything usually fails at everything. "
     "Narrow, specialized agents that each own a single task produce more precise, reliable results."),
    ("2", ACCENT,
     "Retrieval-Augmented Generation (RAG):  ",
     "Instead of relying on general training data, RAG agents actively query a live vector database "
     "of the real codebase — giving them precise, up-to-date architectural knowledge before they respond."),
    ("3", GREEN,
     "Iterative Debate:  ",
     "Agents challenge each other in a structured loop. The Review Agent critiques the Refactor Agent's "
     "suggestion; they iterate until the output meets a defined quality threshold or a timeout fires."),
]

yr = 1.5
for num, col, bold_part, rest in rows:
    box(s3, 0.32, yr, 12.7, 1.58, fill=RGBColor(0x0A,0x1A,0x30))
    box(s3, 0.35, yr+0.32, 0.55, 0.55, fill=col)
    tb(s3, num, 0.35, yr+0.30, 0.55, 0.55, sz=16, bold=True,
       align=PP_ALIGN.CENTER)
    two_run_line(s3, 1.05, yr+0.1, 11.5, 1.35,
                 bold_part, ACCENT, 15, True,
                 rest, LIGHT_GRAY, 14)
    yr += 1.76

pgnum(s3, 3)

# ─── SLIDE 4 — Architecture ─────────────────────

s4 = new_slide()
header(s4, "Proposed Conceptual Architecture", "System Design & Data Flow")

flow_boxes = [
    ("GitHub\nWebhook",               ACCENT2),
    ("PR Diff\nExtractor",            RGBColor(0x15,0x3A,0x6A)),
    ("Vector DB\n(Pinecone)",         RGBColor(0x15,0x3A,0x6A)),
    ("State-Graph\nOrchestrator\n(LangGraph)", ACCENT),
    ("Agent\nSwarm",                  GREEN),
]

bx, by, bw, bh, gap = 0.32, 2.3, 2.18, 1.15, 0.3

for i, (label, col) in enumerate(flow_boxes):
    box(s4, bx, by, bw, bh, fill=col)
    tb(s4, label, bx, by+0.1, bw, bh-0.1,
       sz=13, bold=True, align=PP_ALIGN.CENTER)
    if i < len(flow_boxes)-1:
        box(s4, bx+bw, by+bh/2-0.045, gap, 0.09, fill=ACCENT)
        box(s4, bx+bw+gap-0.04, by+bh/2-0.13, 0.14, 0.29, fill=ACCENT)
    bx += bw + gap

# Debate loop arrow
box(s4, 8.85, 3.63, 4.12, 0.07, fill=ACCENT2)
box(s4, 8.85, 3.63, 0.07, 0.78, fill=ACCENT2)
tb(s4, "↺  Iterative Debate / Consensus Loop",
   9.0, 3.72, 5.0, 0.38, sz=12, color=ACCENT2, italic=True)

notes = [
    "▸  A GitHub webhook fires on every PR merge — the swarm intercepts it instantly.",
    "▸  The PR diff is extracted and converted to vector embeddings for semantic search.",
    "▸  Pinecone is queried for historically coupled modules and related commit history.",
    "▸  LangGraph orchestrates cyclic, stateful handoffs between agents — no single point of failure.",
]
yn = 4.55
for note in notes:
    tb(s4, note, 0.4, yn, 12.5, 0.38, sz=13, color=LIGHT_GRAY)
    yn += 0.43

pgnum(s4, 4)

# ─── SLIDE 5 — Shadow Team ──────────────────────

s5 = new_slide()
header(s5, 'The "Shadow Team" Workflow', "Agent Roles in Action")

agents = [
    (ACCENT2, "🔍  Triage Agent",
     "Watches the GitHub webhook stream. Categorizes each incoming PR (bug fix, feature, refactor, "
     "security patch) and routes it to the correct downstream agents with appropriate priority."),
    (ACCENT,  "📦  Retrieval Agent",
     "Queries the vector database for all files historically coupled to the changed code, past commit "
     "messages, and related issue threads — bringing full context into the swarm's shared memory."),
    (YELLOW,  "🧠  Review Agent",
     "Reads the PR diff alongside retrieved context and flags logic errors, broken interfaces, and "
     "architectural anti-patterns. Acts as the swarm's internal critic and quality gatekeeper."),
    (GREEN,   "🔧  Refactor Agent",
     "Proposes the corrected code. Enters a structured debate loop with the Review Agent, iterating "
     "until the Review Agent approves or a maximum round limit is reached."),
]

ya = 1.35
for col, title, body in agents:
    box(s5, 0.28, ya, 12.75, 1.36, fill=DARK_CARD)
    box(s5, 0.28, ya, 0.12, 1.36, fill=col)
    tb(s5, title, 0.52, ya+0.05, 4.5, 0.45, sz=15, bold=True, color=col)
    tb(s5, body,  0.52, ya+0.5, 12.2, 0.78, sz=13, color=LIGHT_GRAY)
    ya += 1.5

pgnum(s5, 5)

# ─── SLIDE 6 — Ethics ──────────────────────────

s6 = new_slide()
header(s6, "Ethical & Social Considerations", "Human & Governance Impact")

ethics = [
    (RED_SOFT, "⚖️  Accountability Dilemma",
     "If an autonomous swarm approves and merges a PR that introduces a security vulnerability, "
     "who is legally responsible — the developer who triggered it, the team lead, or the AI vendor?"),
    (YELLOW,   "👩‍💻  The Junior Developer Squeeze",
     "PR reviews are how junior engineers learn to read other people's code in depth. Automating them "
     "risks producing developers who never develop that skill — weakening teams long-term."),
    (ACCENT,   "🔒  Security & Bias Risks",
     "Agents may hallucinate secure-looking but vulnerable dependency versions, or enforce biased "
     "coding patterns derived from skewed training corpora — silently spreading bad practices at scale."),
]

ye = 1.38
for col, title, body in ethics:
    box(s6, 0.28, ye, 12.75, 1.72, fill=DARK_CARD)
    box(s6, 0.28, ye, 12.75, 0.06, fill=col)
    tb(s6, title, 0.48, ye+0.15, 12.1, 0.45, sz=15, bold=True, color=col)
    tb(s6, body,  0.48, ye+0.6,  12.1, 1.0,  sz=13.5, color=LIGHT_GRAY)
    ye += 1.9

pgnum(s6, 6)

# ─── SLIDE 7 — Limitations & Future ────────────

s7 = new_slide()
header(s7, "Limitations & Future Scope", "A Realistic, Forward-Looking View")

# Left panel
box(s7, 0.28, 1.35, 6.1, 4.65, fill=RGBColor(0x10,0x20,0x38))
box(s7, 0.28, 1.35, 6.1, 0.06, fill=RED_SOFT)
tb(s7, "⚠️  Current Constraints",
   0.44, 1.46, 5.8, 0.45, sz=15, bold=True, color=RGBColor(0xFF,0x88,0x88))

limits = [
    "Embedding large enterprise monorepos is expensive — token limits are hit quickly and costs spike.",
    "Agent debate loops can turn infinite without a well-tuned timeout and fallback exit condition.",
    "Vector search may retrieve irrelevant code context, sending agents down the wrong analytical path.",
]
yl = 2.02
for lim in limits:
    tb(s7, "▸  " + lim, 0.44, yl, 5.82, 0.95, sz=12.5, color=LIGHT_GRAY)
    yl += 1.08

# Right panel
box(s7, 6.95, 1.35, 6.1, 4.65, fill=RGBColor(0x0A,0x25,0x1A))
box(s7, 6.95, 1.35, 6.1, 0.06, fill=GREEN)
tb(s7, "🚀  Future Scope",
   7.1, 1.46, 5.8, 0.45, sz=15, bold=True, color=GREEN)

future = [
    "Execution Agent: Spins up a secure sandbox to compile and empirically test the fix before merging.",
    "Self-Healing Pipelines: The swarm monitors production logs and autonomously opens PRs to fix live regressions.",
    "Narrator Agent: Generates a human-readable audit trail explaining every decision the swarm made.",
]
yf = 2.02
for fut in future:
    tb(s7, "▸  " + fut, 7.1, yf, 5.82, 0.95, sz=12.5, color=LIGHT_GRAY)
    yf += 1.08

# Conclusion banner
box(s7, 0.28, 6.2, 12.77, 0.9, fill=RGBColor(0x00,0x30,0x50))
tb(s7,
   "Conclusion:  Autonomous Multi-Agent Swarms shift software engineering from manual validation "
   "to automated self-healing — making codebases faster, more reliable, and continuously improving.",
   0.48, 6.26, 12.3, 0.8, sz=13, align=PP_ALIGN.CENTER, italic=True)

pgnum(s7, 7)

# ─── Save ───────────────────────────────────────
output = r"D:\Major project\Multi_Agent_Swarm_PPT.pptx"
prs.save(output)
print("Saved →", output)
